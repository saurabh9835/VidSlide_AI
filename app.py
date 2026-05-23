import os
import shutil
import time
import uuid
import zipfile
import logging

from flask import (
    Flask,
    render_template,
    request,
    redirect,
    url_for,
    session,
    send_file,
    send_from_directory,
    flash,
    after_this_request,
    jsonify,
)
import traceback
import importlib

# Optional ML imports (heavy)
import cv2
import numpy as np
import yt_dlp
from pptx import Presentation
from pptx.util import Inches
from sklearn.metrics.pairwise import cosine_similarity

ResNet50 = None
preprocess_input = None
keras_image = None

def load_keras_modules():
    global ResNet50, preprocess_input, keras_image
    try:
        resnet_mod = importlib.import_module("tensorflow.keras.applications.resnet50")
        keras_image = importlib.import_module("tensorflow.keras.preprocessing.image")
    except ModuleNotFoundError:
        resnet_mod = importlib.import_module("keras.applications.resnet50")
        keras_image = importlib.import_module("keras.preprocessing.image")
    ResNet50 = getattr(resnet_mod, "ResNet50")
    preprocess_input = getattr(resnet_mod, "preprocess_input")

# -------------------------
# Configuration
# -------------------------
app = Flask(__name__, static_folder="static")
app.secret_key = os.environ.get("FLASK_APP_SECRET_KEY", "dev-secret-key")
logging.basicConfig(level=logging.INFO)

# folders
SESSIONS_DIR = "sessions"
os.makedirs(SESSIONS_DIR, exist_ok=True)

# ML model (ResNet50)
# If you want to skip heavy loading during dev, set SKIP_ML=True environment var.
SKIP_ML = os.environ.get("SKIP_ML", "0") == "1"
if not SKIP_ML:
    try:
        load_keras_modules()
        base_model = ResNet50(weights="imagenet", include_top=False, pooling="avg")
        app.logger.info("ResNet50 loaded.")
    except Exception as e:
        app.logger.error(f"Failed to load ResNet50: {e}")
        base_model = None
else:
    base_model = None
    app.logger.info("Skipping ML model load (SKIP_ML=1).")

# -------------------------
# Utilities
# -------------------------
def cleanup_old_sessions(max_age_seconds=3600):
    """Remove session folders older than max_age_seconds."""
    now = time.time()
    if not os.path.exists(SESSIONS_DIR):
        return

    for sid in os.listdir(SESSIONS_DIR):
        path = os.path.join(SESSIONS_DIR, sid)
        if not os.path.isdir(path):
            continue
        try:
            last_mod = os.path.getmtime(path)
            if now - last_mod > max_age_seconds:
                shutil.rmtree(path, ignore_errors=True)
                app.logger.info(f"[CLEANUP] Removed old session: {sid}")
        except Exception as e:
            app.logger.error(f"Error cleaning session {sid}: {e}")

def download_video_url(video_url, session_dir):
    """Download a remote video URL into the session directory and return local path.

    Using a local file is more reliable for OpenCV VideoCapture than a remote stream URL.
    """
    out_path = os.path.join(session_dir, "video.mp4")
    ydl_opts = {
        "format": "bestvideo+bestaudio/best",
        "outtmpl": out_path,
        "merge_output_format": "mp4",
        "noplaylist": True,
        "quiet": True,
    }
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([video_url])
    except Exception as e:
        msg = str(e)
        # Detect common yt-dlp DRM message
        if 'DRM' in msg or 'DRM protected' in msg or 'This video is DRM protected' in msg:
            raise RuntimeError("Video appears to be DRM-protected and cannot be downloaded.")
        raise RuntimeError(f"yt-dlp download error: {e}")

    if not os.path.exists(out_path):
        raise RuntimeError("Download failed; video file not found")
    return out_path

def extract_features(frame):
    """Given a BGR OpenCV frame, return a ResNet50 feature vector (1D numpy)."""
    if base_model is None:
        # Dummy features fallback (so code works if model not loaded)
        return np.random.random(2048).astype(np.float32)

    img = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
    img = cv2.resize(img, (224, 224))
    arr = keras_image.img_to_array(img)
    arr = np.expand_dims(arr, axis=0)
    arr = preprocess_input(arr)
    feats = base_model.predict(arr, verbose=0)
    return feats.flatten()

def video_url_to_slides(video_url, session_id, similarity_threshold=0.98):
    """
    Process a remote video URL and extract slide images using ML logic with
    user-defined similarity threshold. Returns list of image filenames and pptx path.
    """
    session_dir = os.path.join(SESSIONS_DIR, session_id)
    slides_dir = os.path.join(session_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    pptx_dir = os.path.join(session_dir, "presentations")
    os.makedirs(pptx_dir, exist_ok=True)
    pptx_path = os.path.join(pptx_dir, "extracted_slides.pptx")

    try:
        video_path = download_video_url(video_url, session_dir)
    except Exception as e:
        raise RuntimeError(f"Could not download video: {e}")

    return process_video(video_path, session_id, similarity_threshold)


def process_video(video_path, session_id, similarity_threshold=0.98):
    """Common video -> slides processing for a local video file path.

    Returns (saved_images, pptx_path, zip_path)
    """
    session_dir = os.path.join(SESSIONS_DIR, session_id)
    slides_dir = os.path.join(session_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    pptx_dir = os.path.join(session_dir, "presentations")
    os.makedirs(pptx_dir, exist_ok=True)
    pptx_path = os.path.join(pptx_dir, "extracted_slides.pptx")

    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise RuntimeError("Could not open video file for processing.")

    # ML state variables
    STABILITY_THRESHOLD_FRAMES = 10
    state = "STABLE"
    last_confirmed_features = None
    last_transition_features = None
    potential_new_frame = None
    frames_stable_in_row = 0

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    frame_count = 0
    slide_count = 0
    saved_images = []

    # Read frames
    while True:
        ret, frame = cap.read()
        if not ret:
            break
        frame_count += 1

        # Sample every 3rd frame to speed up
        if frame_count % 3 != 0:
            continue

        current_features = extract_features(frame)

        if last_confirmed_features is None:
            # always accept the first sampled frame as the first slide
            slide_count += 1
            fname = f"slide_{slide_count}.png"
            path = os.path.join(slides_dir, fname)
            cv2.imwrite(path, frame, [cv2.IMWRITE_PNG_COMPRESSION, 3])
            saved_images.append(fname)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            last_confirmed_features = current_features
            continue

        if state == "STABLE":
            sim = cosine_similarity(
                last_confirmed_features.reshape(1, -1),
                current_features.reshape(1, -1)
            )[0][0]
            if sim < similarity_threshold:
                state = "IN_TRANSITION"
                last_transition_features = current_features
                frames_stable_in_row = 0

        elif state == "IN_TRANSITION":
            sim_t = cosine_similarity(
                last_transition_features.reshape(1, -1),
                current_features.reshape(1, -1)
            )[0][0]
            if sim_t > 0.99:
                frames_stable_in_row += 1
            else:
                frames_stable_in_row = 0

            last_transition_features = current_features
            potential_new_frame = frame

            if frames_stable_in_row >= STABILITY_THRESHOLD_FRAMES:
                # Accept new slide
                last_confirmed_features = current_features
                slide_count += 1
                fname = f"slide_{slide_count}.png"
                path = os.path.join(slides_dir, fname)
                cv2.imwrite(path, potential_new_frame, [cv2.IMWRITE_PNG_COMPRESSION, 3])
                saved_images.append(fname)

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

                state = "STABLE"
                frames_stable_in_row = 0

    cap.release()
    prs.save(pptx_path)

    # Create zip for slides
    zip_path = os.path.join(session_dir, "slides.zip")
    with zipfile.ZipFile(zip_path, "w") as zf:
        for img in saved_images:
            zf.write(os.path.join(slides_dir, img), arcname=f"slides/{img}")

    return saved_images, pptx_path, zip_path

# -------------------------
# Routes
# -------------------------
@app.route("/", methods=["GET", "POST"])
def index():
    # cleanup old sessions on homepage load
    cleanup_old_sessions(max_age_seconds=3600)

    # Keep original POST handler for non-AJAX fallback. AJAX route `/api/generate`
    return render_template("index.html")


@app.route('/api/generate', methods=['POST'])
def api_generate():
    video_url = request.form.get("video_url", "").strip()
    sensitivity_raw = request.form.get("sensitivity", "0.98")

    # Validate sensitivity
    try:
        sensitivity = float(sensitivity_raw)
        if not (0.90 <= sensitivity <= 0.999):
            sensitivity = 0.98
    except Exception:
        sensitivity = 0.98

    if not video_url:
        return jsonify({'success': False, 'error': 'Please enter a video URL (YouTube or other supported video URL).'}), 400

    session_id = str(uuid.uuid4())
    session['session_id'] = session_id
    session_dir = os.path.join(SESSIONS_DIR, session_id)
    os.makedirs(session_dir, exist_ok=True)

    try:
        images, pptx_path, zip_path = video_url_to_slides(video_url, session_id, similarity_threshold=sensitivity)
    except Exception as e:
        app.logger.error(f"Error processing video: {e}")
        app.logger.error(traceback.format_exc())
        # cleanup on failure
        shutil.rmtree(session_dir, ignore_errors=True)
        return jsonify({'success': False, 'error': f'Error processing video: {e}'}), 500

    # Save paths in session for preview and downloads
    session['images'] = images
    session['pptx_path'] = pptx_path
    session['zip_path'] = zip_path

    return jsonify({'success': True, 'preview_url': url_for('preview')})


@app.route('/api/upload', methods=['POST'])
def api_upload():
    # Accept a video file upload and process it into slides
    if 'video_file' not in request.files:
        return jsonify({'success': False, 'error': 'No file part (video_file) in request.'}), 400
    file = request.files['video_file']
    if file.filename == '':
        return jsonify({'success': False, 'error': 'No selected file.'}), 400

    sensitivity_raw = request.form.get('sensitivity', '0.98')
    try:
        sensitivity = float(sensitivity_raw)
        if not (0.90 <= sensitivity <= 0.999):
            sensitivity = 0.98
    except Exception:
        sensitivity = 0.98

    session_id = str(uuid.uuid4())
    session['session_id'] = session_id
    session_dir = os.path.join(SESSIONS_DIR, session_id)
    os.makedirs(session_dir, exist_ok=True)

    # Save uploaded file to session_dir
    filename = os.path.join(session_dir, 'uploaded_video')
    # try to preserve extension if present
    orig = file.filename
    _, ext = os.path.splitext(orig)
    if ext:
        filename = filename + ext
    else:
        filename = filename + '.mp4'

    try:
        file.save(filename)
    except Exception as e:
        app.logger.error(f"Failed saving uploaded file: {e}")
        return jsonify({'success': False, 'error': 'Failed to save uploaded file.'}), 500

    try:
        images, pptx_path, zip_path = process_video(filename, session_id, similarity_threshold=sensitivity)
    except Exception as e:
        app.logger.error(f"Error processing uploaded video: {e}")
        app.logger.error(traceback.format_exc())
        shutil.rmtree(session_dir, ignore_errors=True)
        return jsonify({'success': False, 'error': f'Error processing video: {e}'}), 500

    session['images'] = images
    session['pptx_path'] = pptx_path
    session['zip_path'] = zip_path

    return jsonify({'success': True, 'preview_url': url_for('preview')})

    return render_template("index.html")

@app.route("/preview")
def preview():
    session_id = session.get("session_id")
    images = session.get("images", [])
    if not session_id or not images:
        return redirect(url_for("index"))

    # preview images will be served from /sessions/<session_id>/slides/<filename>
    return render_template("preview.html", session_id=session_id, images=images)

@app.route("/sessions/<session_id>/slides/<filename>")
def serve_slide(session_id, filename):
    slide_dir = os.path.join(SESSIONS_DIR, session_id, "slides")
    if not os.path.exists(os.path.join(slide_dir, filename)):
        return "Slide not found", 404
    return send_from_directory(slide_dir, filename)

@app.route("/download_pptx/<session_id>")
def download_pptx(session_id):
    pptx_path = os.path.join(SESSIONS_DIR, session_id, "presentations", "extracted_slides.pptx")
    if not os.path.exists(pptx_path):
        return "PPTX not found", 404
    # OPTIONAL: cleanup after sending (comment out if you want manual cleanup)
    @after_this_request
    def _cleanup(response):
        # don't delete entire session; let global cleanup handle it
        return response
    return send_file(pptx_path, as_attachment=True, download_name="Extracted_Slides.pptx")

@app.route("/download_zip/<session_id>")
def download_zip(session_id):
    zip_path = os.path.join(SESSIONS_DIR, session_id, "slides.zip")
    if not os.path.exists(zip_path):
        return "ZIP not found", 404
    return send_file(zip_path, as_attachment=True, download_name="slides.zip")

# -------------------------
# Run
# -------------------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
